from __future__ import annotations

from pathlib import Path
import shutil

from .engine import RedactionEntry


OFFICE_SUFFIXES = {".docx", ".xlsx", ".pptx"}


def read_office_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".docx":
        return _read_docx_text(path)
    if suffix == ".xlsx":
        return _read_xlsx_text(path)
    if suffix == ".pptx":
        return _read_pptx_text(path)
    raise ValueError(f"Formato Office non supportato: {path.name}")


def write_redacted_office(source_path: Path, output_path: Path, entries: list[RedactionEntry]) -> None:
    _write_office_with_replacements(
        source_path,
        output_path,
        [(entry.value, entry.tag) for entry in entries],
    )


def write_restored_office(redacted_path: Path, output_path: Path, entries: list[dict]) -> None:
    _write_office_with_replacements(
        redacted_path,
        output_path,
        [(str(entry["tag"]), str(entry["value"])) for entry in entries],
    )


def _write_office_with_replacements(
    source_path: Path,
    output_path: Path,
    replacements: list[tuple[str, str]],
) -> None:
    suffix = source_path.suffix.lower()
    output_path.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(source_path, output_path)
    if suffix == ".docx":
        _replace_docx(output_path, replacements)
    elif suffix == ".xlsx":
        _replace_xlsx(output_path, replacements)
    elif suffix == ".pptx":
        _replace_pptx(output_path, replacements)
    else:
        raise ValueError(f"Formato Office non supportato: {source_path.name}")


def _read_docx_text(path: Path) -> str:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("Per leggere DOCX installa python-docx.") from exc

    document = Document(str(path))
    parts: list[str] = []
    parts.extend(paragraph.text for paragraph in document.paragraphs)
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                parts.extend(paragraph.text for paragraph in cell.paragraphs)
    return "\n".join(part for part in parts if part)


def _read_xlsx_text(path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise RuntimeError("Per leggere XLSX installa openpyxl.") from exc

    workbook = load_workbook(path, data_only=False)
    parts: list[str] = []
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                if isinstance(cell.value, str):
                    parts.append(cell.value)
    workbook.close()
    return "\n".join(parts)


def _read_pptx_text(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("Per leggere PPTX installa python-pptx.") from exc

    presentation = Presentation(str(path))
    parts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text:
                            parts.append(cell.text)
    return "\n".join(parts)


def _replace_docx(path: Path, replacements: list[tuple[str, str]]) -> None:
    from docx import Document

    document = Document(str(path))
    for paragraph in document.paragraphs:
        _replace_in_paragraph(paragraph, replacements)
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    _replace_in_paragraph(paragraph, replacements)
    document.save(str(path))


def _replace_xlsx(path: Path, replacements: list[tuple[str, str]]) -> None:
    from openpyxl import load_workbook

    workbook = load_workbook(path)
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                if isinstance(cell.value, str):
                    cell.value = _replace_once(cell.value, replacements)
    workbook.save(path)
    workbook.close()


def _replace_pptx(path: Path, replacements: list[tuple[str, str]]) -> None:
    from pptx import Presentation

    presentation = Presentation(str(path))
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    _replace_in_pptx_paragraph(paragraph, replacements)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        cell.text = _replace_once(cell.text, replacements)
    presentation.save(str(path))


def _replace_in_paragraph(paragraph, replacements: list[tuple[str, str]]) -> None:
    if not paragraph.runs:
        return
    original = paragraph.text
    replaced = _replace_once(original, replacements)
    if replaced == original:
        return
    paragraph.runs[0].text = replaced
    for run in paragraph.runs[1:]:
        run.text = ""


def _replace_in_pptx_paragraph(paragraph, replacements: list[tuple[str, str]]) -> None:
    if not paragraph.runs:
        return
    original = "".join(run.text for run in paragraph.runs)
    replaced = _replace_once(original, replacements)
    if replaced == original:
        return
    paragraph.runs[0].text = replaced
    for run in paragraph.runs[1:]:
        run.text = ""


def _replace_once(text: str, replacements: list[tuple[str, str]]) -> str:
    result = text
    for old, new in replacements:
        if old:
            result = result.replace(old, new, 1)
    return result
